import pandas as pd
from pptx import Presentation

class doc_tools():
    def __init__(self,project_meta,outline):
        self.project_meta,self.outline = project_meta,outline
        self.compiled_outline = self.write_outline()

    def write_outline(self):
        tasks = (self.project_meta['completed_tasks'].iloc[0],self.project_meta['total_tasks'].iloc[0])
        sect_sep = ('_'*25) + '\n\n'
        _title = f"Project: {self.project_meta['project'].iloc[0]}\n\n"
        _lead = f"Lead: {self.project_meta['project_lead'].iloc[0]}\n\n"
        _stat = f"Status: \n\n\t{self.project_meta['project_status'].iloc[0]}\n\n\tTasks: {tasks[0]} of {tasks[1]} complete\n\n\tEst. Completion: {self.project_meta['est_completion'].iloc[0]}\n"
        _desc = f"Summary:\n\t{self.project_meta['project_desc'].iloc[0]}\n"
        _outhd = sect_sep + f'\noutline:\n'
        _out = ''
        for i in range(len(self.outline)):
            row = self.outline.iloc[i]
            out = f"""

            Task ID: {row['task_id']}
            Task Name: {row['task_name']}
            Task Owner: {row['owner']}
            Task Status: {row['task_status']}
            Task Desc:
                {row['task_desc']}
            Task Dependencies: {row['task_dependencies']}
            Tables Affected: {row['tbls_affected']}
            Created: {row['create_date']}, Est Completion: {row['est_completion']}

            """
            _out += out
        compiled_doc = _title + _lead + _stat + _desc + _outhd + _out
        return compiled_doc

    def ci_request(self,pptx_document):
        prs = Presentation(pptx_document)
        text_runs = {}
        for sldx,slide in enumerate(prs.slides):
            text_runs[sldx] = {}
            for sdx,shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                text_runs[sldx][sdx] = []
                for pdx,paragraph in enumerate(shape.text_frame.paragraphs):
                    for rdx,run in enumerate(paragraph.runs):
                        text_runs[sldx][sdx].append(run.text)
        title = text_runs.get(0)
        title = '\n'.join([x for y in title.values() for x in y if x])
        for key,val in text_runs.items():
            for k,v in val.items():
                if 'CI Request' in [x.strip() for x in v]:
                    cir_k = key
                if 'Business Requirements' in [x.strip() for x in v]:
                    brq_k = key
        ci_request = text_runs.get(cir_k)
        cir_body = '\n'.join([x for y in ci_request.values() for x in y if x])
        b_require = text_runs.get(brq_k)
        brq_body = '\n'.join([x for y in b_require.values() for x in y if x])
        request = title + '\n\n' + cir_body + '\n\n' + brq_body + '\n__________\n\n'
        return request